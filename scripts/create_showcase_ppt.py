from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "runs" / "presentation_assets"
ASCII_OUTPUT = ROOT / "project_showcase_presentation.pptx"
CN_OUTPUT = ROOT / "成果展示PPT_完成版.pptx"

BG = RGBColor(245, 245, 247)
CARD = RGBColor(255, 255, 255)
TEXT = RGBColor(29, 29, 31)
MUTED = RGBColor(110, 110, 115)
LINE = RGBColor(224, 224, 230)
ACCENT = RGBColor(0, 113, 227)
ACCENT_SOFT = RGBColor(233, 244, 255)
GREEN = RGBColor(52, 199, 89)
ORANGE = RGBColor(255, 159, 10)
RED = RGBColor(255, 69, 58)

FONT = "Microsoft YaHei"
EN_FONT = "Aptos"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def ensure_assets_dir() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    circle_a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.9), Inches(-0.7), Inches(3.2), Inches(3.2))
    circle_a.fill.solid()
    circle_a.fill.fore_color.rgb = ACCENT
    circle_a.fill.transparency = 0.88
    circle_a.line.fill.background()

    circle_b = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.7), Inches(6.0), Inches(2.8), Inches(2.8))
    circle_b.fill.solid()
    circle_b.fill.fore_color.rgb = GREEN
    circle_b.fill.transparency = 0.9
    circle_b.line.fill.background()


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    font_size: float = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(
    shape,
    lines: Iterable[str],
    font_size: float = 18,
    color: RGBColor = TEXT,
    bullet: bool = False,
    level: int = 0,
    space_after: int = 6,
):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    first = True
    for line in lines:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.text = line
        paragraph.level = level
        paragraph.space_after = Pt(space_after)
        paragraph.bullet = bullet
        run = paragraph.runs[0]
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_card(slide, left: float, top: float, width: float, height: float, fill_color: RGBColor = CARD):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    return shape


def add_chip(slide, left: float, top: float, width: float, height: float, text: str, fill_color: RGBColor = ACCENT_SOFT, text_color: RGBColor = ACCENT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = text_color
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_title(slide, slide_index: int, eyebrow: str, title: str, subtitle: str) -> None:
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(1.1), Inches(0.35), f"{slide_index:02d}", 16, True, ACCENT, PP_ALIGN.CENTER, EN_FONT)
    add_textbox(slide, Inches(1.35), Inches(0.45), Inches(3.0), Inches(0.3), eyebrow, 11, True, ACCENT)
    add_textbox(slide, Inches(0.7), Inches(0.8), Inches(8.4), Inches(0.6), title, 28, True, TEXT)
    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(8.8), Inches(0.45), subtitle, 12, False, MUTED)


def add_footer(slide, page_no: int) -> None:
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(7.0), Inches(11.95), Inches(0.02)).fill.solid()
    footer_bar = slide.shapes[-1]
    footer_bar.fill.fore_color.rgb = LINE
    footer_bar.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(7.05), Inches(5.6), Inches(0.2), "人工智能应用实践 | 乐高积木图像识别项目成果展示", 9, False, MUTED)
    add_textbox(slide, Inches(12.0), Inches(7.03), Inches(0.6), Inches(0.2), str(page_no), 9, False, MUTED, PP_ALIGN.RIGHT, EN_FONT)


def add_stat_card(slide, left: float, top: float, width: float, height: float, label: str, value: str, accent_color: RGBColor = ACCENT):
    add_card(slide, left, top, width, height)
    add_textbox(slide, left + Inches(0.22), top + Inches(0.16), width - Inches(0.44), Inches(0.24), label, 11, True, MUTED)
    add_textbox(slide, left + Inches(0.22), top + Inches(0.45), width - Inches(0.44), Inches(0.45), value, 24, True, accent_color)


def add_picture_contain(slide, path: Path, left: float, top: float, width: float, height: float):
    with Image.open(path) as image:
        image_width, image_height = image.size
    target_ratio = width / height
    image_ratio = image_width / image_height
    if image_ratio > target_ratio:
        scaled_width = width
        scaled_height = width / image_ratio
        top += (height - scaled_height) / 2
    else:
        scaled_height = height
        scaled_width = height * image_ratio
        left += (width - scaled_width) / 2
    slide.shapes.add_picture(str(path), left, top, width=scaled_width, height=scaled_height)


def build_sample_montage() -> Path:
    output_path = ASSET_DIR / "sample_montage.png"
    if output_path.exists():
        return output_path

    root = ROOT / "data" / "splits_stage2_raw" / "test"
    classes = sorted([path.name for path in root.iterdir() if path.is_dir()])
    card_w, card_h = 360, 300
    margin = 36
    cols = 4
    rows = 2
    canvas = Image.new("RGB", (cols * card_w + (cols + 1) * margin, rows * card_h + (rows + 1) * margin), (245, 245, 247))
    draw = ImageDraw.Draw(canvas)

    for index, class_name in enumerate(classes):
        row = index // cols
        col = index % cols
        x = margin + col * (card_w + margin)
        y = margin + row * (card_h + margin)
        draw.rounded_rectangle((x, y, x + card_w, y + card_h), radius=28, fill=(255, 255, 255), outline=(224, 224, 230), width=2)
        image_path = sorted((root / class_name).glob("*"))[0]
        with Image.open(image_path) as image:
            image = image.convert("RGB")
            image.thumbnail((card_w - 44, 190))
            image_x = x + (card_w - image.width) // 2
            image_y = y + 24
            canvas.paste(image, (image_x, image_y))
        draw.text((x + 24, y + 230), class_name, fill=(29, 29, 31))
        draw.text((x + 24, y + 260), "测试集代表样本", fill=(110, 110, 115))

    canvas.save(output_path)
    return output_path


def build_current_class_totals(split_report: dict) -> dict[str, int]:
    totals: dict[str, int] = {}
    for split_counts in split_report.values():
        for class_name, count in split_counts.items():
            totals[class_name] = totals.get(class_name, 0) + int(count)
    return totals


def add_accuracy_chart(slide, left: float, top: float, width: float, height: float, baseline: dict, enhanced: dict):
    categories = ["1x1", "1x2", "1x3", "1x4", "2x2", "2x3", "2x4"]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("基础推理", [round(float(baseline["per_class_accuracy"][item]) * 100, 1) for item in categories])
    chart_data.add_series("抗阴影双路融合", [round(float(enhanced["per_class_accuracy"][item]) * 100, 1) for item in categories])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 20
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.name = EN_FONT
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = EN_FONT
    chart.category_axis.format.line.color.rgb = LINE
    chart.value_axis.format.line.color.rgb = LINE

    for series, color in zip(chart.series, (ACCENT, GREEN)):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color
        series.has_data_labels = True
        series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        series.data_labels.font.size = Pt(9)
        series.data_labels.font.name = EN_FONT


def add_table(slide, left: float, top: float, width: float, height: float, headers: list[str], rows: list[list[str]]):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for index, header in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_SOFT
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = ACCENT
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT if col_index == 0 else EN_FONT
                    run.font.size = Pt(11)
                    run.font.color.rgb = TEXT
    return table


def build_presentation() -> Presentation:
    ensure_assets_dir()

    stage1_summary = load_json(ROOT / "runs" / "stage3" / "verification_summary.json")
    current_split = load_json(ROOT / "data" / "splits_stage2_raw" / "split_report.json")
    baseline_metrics = load_json(ROOT / "runs" / "stage2" / "improve_raw128_autocrop_lr5e4_7class_v6" / "test_metrics_recheck_raw_20260422.json")
    enhanced_metrics = load_json(ROOT / "runs" / "stage2" / "shadow_tta_v6_20260422" / "codepath_recheck.json")
    model_summary = load_json(ROOT / "runs" / "stage2" / "improve_raw128_autocrop_lr5e4_7class_v6" / "model_summary.json")
    depth_width_summary = load_json(ROOT / "runs" / "stage2" / "experiments" / "depth_width" / "summary.json")
    batch_size_summary = load_json(ROOT / "runs" / "stage2" / "experiments" / "batch_size" / "summary.json")
    dropout_summary = load_json(ROOT / "runs" / "stage2" / "experiments" / "dropout" / "summary.json")

    sample_montage = build_sample_montage()
    history_accuracy = ROOT / "runs" / "stage2" / "improve_raw128_autocrop_lr5e4_7class_v6" / "history_accuracy.png"
    history_loss = ROOT / "runs" / "stage2" / "improve_raw128_autocrop_lr5e4_7class_v6" / "history_loss.png"
    depth_width_chart = ROOT / "runs" / "stage2" / "experiments" / "depth_width" / "summary_test_accuracy.png"
    batch_size_chart = ROOT / "runs" / "stage2" / "experiments" / "batch_size" / "summary_test_accuracy.png"

    current_totals = build_current_class_totals(current_split)
    current_total_count = sum(current_totals.values())
    base_total_count = int(stage1_summary["stage1"]["total_images"])
    test_total = int(enhanced_metrics["total"])
    model_params = int(model_summary["parameter_count"])
    model_config = model_summary["config"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_card(slide, Inches(0.7), Inches(0.7), Inches(7.4), Inches(4.95))
    add_textbox(slide, Inches(1.0), Inches(1.05), Inches(1.5), Inches(0.3), "人工智能应用实践", 12, True, ACCENT)
    add_textbox(slide, Inches(1.0), Inches(1.45), Inches(6.4), Inches(1.2), "乐高积木图像识别项目\n成果展示", 26, True, TEXT)
    add_textbox(slide, Inches(1.0), Inches(2.75), Inches(6.2), Inches(0.8), "基于 PyTorch + Streamlit + Plotly，完成从数据准备、CNN 训练到网页部署与鲁棒性优化的完整链路。", 14, False, MUTED)
    add_chip(slide, Inches(1.0), Inches(3.65), Inches(1.2), Inches(0.34), "PyTorch")
    add_chip(slide, Inches(2.3), Inches(3.65), Inches(1.2), Inches(0.34), "Streamlit")
    add_chip(slide, Inches(3.6), Inches(3.65), Inches(1.0), Inches(0.34), "Plotly")
    add_chip(slide, Inches(4.75), Inches(3.65), Inches(1.1), Inches(0.34), "GPU 训练")
    add_textbox(slide, Inches(1.0), Inches(4.45), Inches(2.6), Inches(0.3), "汇报主题：课程阶段成果展示", 11, False, MUTED)
    add_textbox(slide, Inches(1.0), Inches(4.75), Inches(2.8), Inches(0.3), "日期：2026-04-22", 11, False, MUTED, font_name=EN_FONT)
    add_stat_card(slide, Inches(8.45), Inches(1.05), Inches(1.8), Inches(1.25), "基础复核准确率", "87.5%")
    add_stat_card(slide, Inches(10.45), Inches(1.05), Inches(1.8), Inches(1.25), "优化后准确率", "93.75%", GREEN)
    add_stat_card(slide, Inches(8.45), Inches(2.55), Inches(1.8), Inches(1.25), "识别类别", "7 类", ORANGE)
    add_stat_card(slide, Inches(10.45), Inches(2.55), Inches(1.8), Inches(1.25), "当前测试集", "32 张", RED)
    add_card(slide, Inches(8.45), Inches(4.05), Inches(3.8), Inches(1.6), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(8.72), Inches(4.28), Inches(3.3), Inches(0.28), "成果关键词", 12, True, ACCENT)
    add_textbox(slide, Inches(8.72), Inches(4.63), Inches(3.2), Inches(0.72), "4 层 CNN、Top-K 概率图、多页面应用、拍照识别、抗阴影双路融合", 13, False, TEXT)
    add_footer(slide, 1)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 2, "PROJECT OVERVIEW", "项目背景与目标", "围绕课程实验要求，构建一个可运行、可展示、可优化的乐高积木图像识别系统。")
    left_card = add_card(slide, Inches(0.7), Inches(1.95), Inches(5.8), Inches(4.45))
    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(2.4), Inches(0.3), "项目背景", 18, True, TEXT)
    add_paragraphs(
        left_card,
        [
            "识别对象为 7 类常见 Brick 零件，兼顾网络图片与实拍图片场景。",
            "目标不是只做离线训练，而是完成从数据准备、模型训练到 Web 应用部署的完整闭环。",
            "在最终展示阶段，重点关注模型可解释性、交互体验和实拍鲁棒性。",
        ],
        font_size=15,
        color=TEXT,
        bullet=True,
        space_after=10,
    )
    goal_positions = [
        (Inches(6.8), Inches(2.05), "目标 1", "搭建 Streamlit 多页面应用，形成清晰的数据应用界面。"),
        (Inches(9.7), Inches(2.05), "目标 2", "完成模型加载、预处理、前向推理、Softmax 与结果解析。"),
        (Inches(6.8), Inches(4.0), "目标 3", "使用 Plotly 展示概率分布，提高结果可读性和专业感。"),
        (Inches(9.7), Inches(4.0), "目标 4", "针对实拍阴影问题做鲁棒性优化，提升落地效果。"),
    ]
    for left, top, label, body in goal_positions:
        add_card(slide, left, top, Inches(2.55), Inches(1.65), fill_color=CARD)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.16), Inches(0.9), Inches(0.2), label, 11, True, ACCENT)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.48), Inches(2.15), Inches(0.9), body, 13, False, TEXT)
    add_textbox(slide, Inches(0.7), Inches(6.55), Inches(11.8), Inches(0.3), "实验环境：Python、PyTorch、Streamlit、Plotly、Pillow、Trae IDE、浏览器、GPU 训练环境。", 11, False, MUTED)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 3, "PIPELINE", "总体技术路线", "按照“数据准备 -> 模型训练 -> 应用开发与优化”三阶段推进，逐步完成课程成果。")
    stage_specs = [
        (Inches(0.85), "阶段一", "项目规划与数据准备", ["目录与环境初始化", "样本清洗、命名与尺寸统一", "类别整理与数据划分"]),
        (Inches(4.55), "阶段二", "模型设计与训练", ["CNN 模型搭建", "超参数对比实验", "最佳权重保存与复核"]),
        (Inches(8.25), "阶段三", "应用开发与优化", ["Streamlit 多页面应用", "Plotly 概率可视化", "抗阴影双路融合"]),
    ]
    for index, (left, tag, title, items) in enumerate(stage_specs):
        add_card(slide, left, Inches(2.1), Inches(3.0), Inches(3.5))
        add_chip(slide, left + Inches(0.22), Inches(2.28), Inches(0.95), Inches(0.32), tag)
        add_textbox(slide, left + Inches(0.22), Inches(2.68), Inches(2.4), Inches(0.32), title, 18, True, TEXT)
        for item_index, item in enumerate(items):
            add_textbox(slide, left + Inches(0.28), Inches(3.18 + item_index * 0.56), Inches(2.45), Inches(0.3), f"• {item}", 13, False, TEXT)
        if index < len(stage_specs) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(3.05), Inches(3.85), left + Inches(3.55), Inches(3.85))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(2.5)
    add_textbox(slide, Inches(0.9), Inches(6.3), Inches(11.6), Inches(0.35), "最终形成了“可训练、可验证、可部署、可展示”的课程项目成果。", 14, True, ACCENT)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 4, "DATASET", "数据集构建与识别任务", "以 7 类 Brick 识别任务为核心，并在项目后期对弱类和实拍样本继续补充整理。")
    add_card(slide, Inches(0.7), Inches(1.95), Inches(6.0), Inches(4.75))
    add_picture_contain(slide, sample_montage, Inches(0.92), Inches(2.15), Inches(5.55), Inches(4.15))
    add_textbox(slide, Inches(7.0), Inches(2.0), Inches(2.5), Inches(0.32), "关键数据规模", 18, True, TEXT)
    add_stat_card(slide, Inches(7.0), Inches(2.35), Inches(1.65), Inches(1.0), "基础整理", f"{base_total_count} 张", ACCENT)
    add_stat_card(slide, Inches(8.82), Inches(2.35), Inches(1.65), Inches(1.0), "当前训练划分", f"{current_total_count} 张", GREEN)
    add_stat_card(slide, Inches(10.65), Inches(2.35), Inches(1.65), Inches(1.0), "测试集复核", f"{test_total} 张", ORANGE)
    class_rows = [[class_name, str(current_totals[class_name])] for class_name in sorted(current_totals)]
    add_table(slide, Inches(7.0), Inches(3.65), Inches(5.3), Inches(2.45), ["类别", "当前样本数"], class_rows)
    add_textbox(slide, Inches(7.0), Inches(6.25), Inches(5.1), Inches(0.35), "说明：基础样本先整理为 213 张，后续针对弱类与实拍场景继续补充，形成当前训练划分。", 10, False, MUTED)
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 5, "MODEL", "CNN 模型结构设计", "主力模型采用 4 层卷积池化循环的 baseline CNN，并配合全局平均池化与全连接输出层。")
    architecture_labels = [
        "输入\n128×128×3",
        "Conv-BN-ReLU-Pool\n48 ch",
        "Conv-BN-ReLU-Pool\n96 ch",
        "Conv-BN-ReLU-Pool\n192 ch",
        "Conv-BN-ReLU-Pool\n384 ch",
        "GAP\n1×1",
        "Dropout + FC",
        "7 类输出",
    ]
    block_left = 0.78
    block_width = 1.35
    for index, label in enumerate(architecture_labels):
        current_left = Inches(block_left + index * 1.53)
        fill_color = ACCENT_SOFT if index in (1, 2, 3, 4) else CARD
        add_card(slide, current_left, Inches(2.55), Inches(block_width), Inches(1.35), fill_color=fill_color)
        add_textbox(slide, current_left + Inches(0.08), Inches(2.87), Inches(block_width - 0.16), Inches(0.72), label, 13, True, TEXT, PP_ALIGN.CENTER)
        if index < len(architecture_labels) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, current_left + Inches(block_width), Inches(3.22), current_left + Inches(1.47), Inches(3.22))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(2.2)
    add_card(slide, Inches(0.9), Inches(4.45), Inches(5.5), Inches(1.65))
    add_textbox(slide, Inches(1.15), Inches(4.68), Inches(2.1), Inches(0.28), "最终配置", 15, True, TEXT)
    add_textbox(
        slide,
        Inches(1.15),
        Inches(5.05),
        Inches(5.0),
        Inches(0.72),
        f"image_size={model_config['image_size']}  batch_size={model_config['batch_size']}  lr={model_config['learning_rate']}  depth={model_config['depth']}  base_channels={model_config['base_channels']}  dropout={model_config['dropout']}  auto_crop={model_config['auto_crop']}",
        12,
        False,
        MUTED,
        font_name=EN_FONT,
    )
    add_card(slide, Inches(6.7), Inches(4.45), Inches(5.55), Inches(1.65))
    add_textbox(slide, Inches(6.95), Inches(4.68), Inches(2.6), Inches(0.28), "模型要点", 15, True, TEXT)
    add_textbox(slide, Inches(6.95), Inches(5.02), Inches(5.0), Inches(0.78), f"参数量约 {model_params / 10000:.1f} 万；输出维度 7 类。\n结构简洁、训练稳定，更适合当前中小规模课程数据集。", 12, False, MUTED)
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 6, "EXPERIMENTS", "超参数实验与模型选择", "通过多组对比实验筛选更稳定的配置，最终收敛到适合本数据规模的 4 层 baseline CNN。")
    add_card(slide, Inches(0.7), Inches(1.95), Inches(6.0), Inches(3.55))
    add_card(slide, Inches(6.95), Inches(1.95), Inches(5.68), Inches(3.55))
    add_picture_contain(slide, depth_width_chart, Inches(0.95), Inches(2.25), Inches(5.5), Inches(2.9))
    add_picture_contain(slide, batch_size_chart, Inches(7.18), Inches(2.25), Inches(5.2), Inches(2.9))
    depth_best = next(item for item in depth_width_summary if item["override"]["depth"] == 4)
    batch_best = max(batch_size_summary, key=lambda item: item["metrics"]["test_accuracy"])
    dropout_best = max(dropout_summary, key=lambda item: item["metrics"]["test_accuracy"])
    add_card(slide, Inches(0.7), Inches(5.8), Inches(11.93), Inches(0.95), fill_color=ACCENT_SOFT)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(6.05),
        Inches(11.2),
        Inches(0.4),
        f"结论：depth={depth_best['override']['depth']}、base_channels={depth_best['override']['base_channels']} 的结构在当前对比组中表现最好；batch_size={batch_best['override']['batch_size']} 最稳；dropout={dropout_best['override']['dropout']} 更适合当前数据规模。",
        13,
        True,
        ACCENT,
    )
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 7, "RESULTS", "主力模型结果与训练过程", "在原始测试集复核中，主力 v6 模型达到 87.5% 准确率，并呈现较稳定的训练收敛过程。")
    total_test_count = sum(int(value) for value in baseline_metrics["per_class_total"].values())
    add_stat_card(slide, Inches(0.75), Inches(1.95), Inches(2.15), Inches(1.25), "原始测试集复核准确率", "87.5%")
    add_stat_card(slide, Inches(3.05), Inches(1.95), Inches(2.15), Inches(1.25), "测试集样本", f"{total_test_count} 张", ORANGE)
    add_stat_card(slide, Inches(5.35), Inches(1.95), Inches(2.15), Inches(1.25), "模型参数量", f"{model_params / 10000:.1f} 万", GREEN)
    add_card(slide, Inches(7.78), Inches(1.95), Inches(4.85), Inches(1.25), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(8.02), Inches(2.18), Inches(4.4), Inches(0.6), "Top1 预测能较好区分 1x1、1x2、1x4、2x2、2x4 等类别，难点主要集中在外形相近和受阴影影响的实拍样本。", 12, False, TEXT)
    add_card(slide, Inches(0.7), Inches(3.5), Inches(5.95), Inches(2.65))
    add_card(slide, Inches(6.98), Inches(3.5), Inches(5.65), Inches(2.65))
    add_picture_contain(slide, history_accuracy, Inches(0.95), Inches(3.8), Inches(5.45), Inches(2.1))
    add_picture_contain(slide, history_loss, Inches(7.18), Inches(3.8), Inches(5.2), Inches(2.1))
    add_footer(slide, 7)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 8, "APPLICATION", "Streamlit 应用功能展示", "应用侧完成了多页面导航、上传识别、拍照识别、Plotly 图表和识别历史等完整体验。")
    app_cards = [
        ("Classification", "上传图片、摄像头拍照、Top-K 概率、识别历史"),
        ("Examples", "动态读取示例图片，一键识别并展示概率分布图"),
        ("About", "汇总三阶段成果、运行方式与关键说明"),
    ]
    for index, (title, body) in enumerate(app_cards):
        left = Inches(0.85 + index * 4.05)
        add_card(slide, left, Inches(2.05), Inches(3.5), Inches(2.05))
        add_textbox(slide, left + Inches(0.22), Inches(2.34), Inches(2.1), Inches(0.3), title, 20, True, TEXT, font_name=EN_FONT)
        add_textbox(slide, left + Inches(0.22), Inches(2.78), Inches(2.95), Inches(0.9), body, 13, False, MUTED)
    add_card(slide, Inches(0.85), Inches(4.45), Inches(3.4), Inches(1.45), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(1.08), Inches(4.72), Inches(2.8), Inches(0.22), "界面能力", 13, True, ACCENT)
    add_textbox(slide, Inches(1.08), Inches(5.03), Inches(2.9), Inches(0.55), "浅色 CSS 美化、双栏布局、空状态提示与异常处理。", 12, False, TEXT)
    add_card(slide, Inches(4.55), Inches(4.45), Inches(3.4), Inches(1.45), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(4.78), Inches(4.72), Inches(2.8), Inches(0.22), "课堂展示能力", 13, True, ACCENT)
    add_textbox(slide, Inches(4.78), Inches(5.03), Inches(2.9), Inches(0.55), "支持现场拍照识别，结果可视化清晰，适合演示。", 12, False, TEXT)
    add_card(slide, Inches(8.25), Inches(4.45), Inches(4.35), Inches(1.45), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(8.48), Inches(4.72), Inches(3.6), Inches(0.22), "结果输出", 13, True, ACCENT)
    add_textbox(slide, Inches(8.48), Inches(5.03), Inches(3.7), Inches(0.55), "预测类别、置信度、Plotly 概率图、识别历史统一在前端展示。", 12, False, TEXT)
    add_footer(slide, 8)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 9, "INFERENCE FLOW", "端到端推理链路与可视化", "将训练好的模型嵌入 Web 应用后，形成清晰的端到端图片识别流程。")
    flow_labels = [
        "上传 / 拍照",
        "图片预处理",
        "模型加载",
        "前向推理",
        "Softmax",
        "结果解析",
        "Plotly 概率图",
    ]
    for index, label in enumerate(flow_labels):
        left = Inches(0.8 + index * 1.77)
        fill_color = ACCENT_SOFT if index in (1, 4, 6) else CARD
        add_card(slide, left, Inches(2.55), Inches(1.38), Inches(1.1), fill_color=fill_color)
        add_textbox(slide, left + Inches(0.08), Inches(2.86), Inches(1.22), Inches(0.55), label, 13, True, TEXT, PP_ALIGN.CENTER)
        if index < len(flow_labels) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(1.38), Inches(3.1), left + Inches(1.66), Inches(3.1))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(2)
    flow_note = add_card(slide, Inches(0.9), Inches(4.35), Inches(5.3), Inches(1.65))
    add_textbox(slide, Inches(1.15), Inches(4.62), Inches(2.6), Inches(0.24), "工程实现要点", 15, True, TEXT)
    add_textbox(slide, Inches(1.15), Inches(4.98), Inches(4.8), Inches(0.76), "将模型推理、图片读取、Top-K 结果组织和前端渲染拆分成独立模块，便于维护和后续替换模型。", 12, False, MUTED)
    safe_note = add_card(slide, Inches(6.45), Inches(4.35), Inches(5.75), Inches(1.65))
    add_textbox(slide, Inches(6.7), Inches(4.62), Inches(2.8), Inches(0.24), "健壮性处理", 15, True, TEXT)
    add_textbox(slide, Inches(6.7), Inches(4.98), Inches(5.1), Inches(0.76), "加入文件格式校验、空状态提示、异常捕获和识别历史记录，保证应用在课堂演示时更稳定。", 12, False, MUTED)
    add_footer(slide, 9)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 10, "ROBUSTNESS", "实拍阴影鲁棒性优化", "针对实拍图片中常见的阴影、背景和光照变化问题，引入抗阴影双路融合策略。")
    add_card(slide, Inches(0.95), Inches(2.35), Inches(2.2), Inches(2.4))
    add_card(slide, Inches(4.2), Inches(2.35), Inches(2.2), Inches(2.4))
    add_card(slide, Inches(7.55), Inches(2.35), Inches(2.2), Inches(2.4), fill_color=ACCENT_SOFT)
    add_card(slide, Inches(10.55), Inches(2.35), Inches(1.85), Inches(2.4))
    add_textbox(slide, Inches(1.2), Inches(2.72), Inches(1.7), Inches(0.5), "原图推理", 20, True, TEXT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.1), Inches(3.35), Inches(1.9), Inches(0.8), "保留原始纹理与颜色信息。", 13, False, MUTED, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.45), Inches(2.72), Inches(1.7), Inches(0.5), "抗阴影预处理", 20, True, TEXT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.35), Inches(3.35), Inches(1.9), Inches(0.8), "削弱阴影和局部亮度差异。", 13, False, MUTED, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.82), Inches(2.72), Inches(1.6), Inches(0.5), "50 / 50 融合", 20, True, ACCENT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.72), Inches(3.35), Inches(1.9), Inches(0.8), "综合两路预测得分，提高稳定性。", 13, False, TEXT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.78), Inches(2.72), Inches(1.35), Inches(0.5), "最终结果", 20, True, TEXT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.72), Inches(3.35), Inches(1.45), Inches(0.8), "输出更稳定的 Top1 / TopK 结果。", 13, False, MUTED, PP_ALIGN.CENTER)
    for start_x, end_x in ((3.15, 4.0), (6.4, 7.35), (9.75, 10.35)):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start_x), Inches(3.55), Inches(end_x), Inches(3.55))
        connector.line.color.rgb = ACCENT
        connector.line.width = Pt(3)
    add_textbox(slide, Inches(0.95), Inches(5.5), Inches(11.2), Inches(0.4), "优化核心不是盲目重训模型，而是结合问题特征改造推理链路，提升真实场景下的识别鲁棒性。", 13, True, ACCENT)
    add_footer(slide, 10)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 11, "IMPROVEMENT", "优化效果对比", "抗阴影双路融合将应用侧复核准确率从 87.5% 提升到 93.75%，并显著改善弱项类别表现。")
    add_stat_card(slide, Inches(0.8), Inches(1.95), Inches(2.15), Inches(1.2), "基础推理", "87.5%")
    add_stat_card(slide, Inches(3.25), Inches(1.95), Inches(2.15), Inches(1.2), "优化后", "93.75%", GREEN)
    add_card(slide, Inches(5.7), Inches(1.95), Inches(2.0), Inches(1.2), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(5.95), Inches(2.22), Inches(1.5), Inches(0.5), f"30 / {enhanced_metrics['total']}", 24, True, ACCENT, PP_ALIGN.CENTER)
    add_card(slide, Inches(8.0), Inches(1.95), Inches(4.45), Inches(1.2))
    add_textbox(slide, Inches(8.25), Inches(2.2), Inches(3.95), Inches(0.52), "1x4、2x2、2x3、2x4 在优化后都达到 100% 类别准确率。", 12, False, TEXT)
    add_card(slide, Inches(0.8), Inches(3.45), Inches(11.65), Inches(2.8))
    add_accuracy_chart(slide, Inches(1.1), Inches(3.72), Inches(11.0), Inches(2.25), baseline_metrics, enhanced_metrics)
    add_footer(slide, 11)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, 12, "SUMMARY", "项目成果总结与后续展望", "课程要求中的核心链路已经完成，当前成果具备训练、部署和课堂展示能力。")
    completed_card = add_card(slide, Inches(0.8), Inches(2.0), Inches(5.65), Inches(3.8))
    add_textbox(slide, Inches(1.05), Inches(2.28), Inches(2.6), Inches(0.28), "本次已完成", 18, True, TEXT)
    add_paragraphs(
        completed_card,
        [
            "完成 7 类 Brick 数据整理、样本补充与训练划分。",
            "完成 4 层 baseline CNN 训练、对比实验与最佳权重复核。",
            "完成 Streamlit 多页面应用、拍照识别、Plotly 概率图和 CSS 美化。",
            "完成针对实拍阴影问题的推理链路优化，准确率提升到 93.75%。",
        ],
        font_size=14,
        color=TEXT,
        bullet=True,
        space_after=10,
    )
    future_card = add_card(slide, Inches(6.75), Inches(2.0), Inches(5.75), Inches(3.8), fill_color=ACCENT_SOFT)
    add_textbox(slide, Inches(7.0), Inches(2.28), Inches(2.8), Inches(0.28), "后续可继续完善", 18, True, ACCENT)
    add_paragraphs(
        future_card,
        [
            "补充更规范的 Plate 类数据，继续推进 Brick + Plate 联合模型训练。",
            "扩大实拍数据规模，进一步验证多角度、多背景下的稳定性。",
            "继续打磨摄像头交互、移动端体验与批量识别等扩展功能。",
        ],
        font_size=14,
        color=TEXT,
        bullet=True,
        space_after=10,
    )
    add_textbox(slide, Inches(0.85), Inches(6.25), Inches(11.3), Inches(0.4), "结论：项目已形成“数据准备 -> 模型训练 -> 应用实现 -> 鲁棒性优化”的完整课程成果，可直接用于结课展示。", 15, True, ACCENT)
    add_footer(slide, 12)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(ASCII_OUTPUT)
    prs.save(CN_OUTPUT)
    print(ASCII_OUTPUT)
    print(CN_OUTPUT)


if __name__ == "__main__":
    main()
