from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable, Sequence

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPT = ROOT / "成果展示PPT_项目优化版_20260423.pptx"
ASSET_DIR = ROOT / "runs" / "presentation_assets" / "20260423"

BRICK_CLASSES = ["1x1", "1x2", "1x3", "1x4", "2x2", "2x3", "2x4"]

BG = RGBColor(247, 244, 237)
SURFACE = RGBColor(255, 255, 255)
TEXT = RGBColor(41, 47, 54)
MUTED = RGBColor(96, 103, 112)
ACCENT = RGBColor(201, 78, 36)
ACCENT_2 = RGBColor(235, 169, 52)
ACCENT_3 = RGBColor(58, 104, 146)
ACCENT_4 = RGBColor(108, 138, 53)
LIGHT_PANEL = RGBColor(236, 231, 222)

TITLE_FONT = "Microsoft YaHei"
EN_FONT = "Bahnschrift"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def compute_real_metrics(report_path: Path) -> dict:
    payload = load_json(report_path)
    samples = list(payload["samples"])
    overall_total = len(samples)
    overall_correct = sum(1 for row in samples if row["correct"])
    original_samples = [row for row in samples if "__rot_" not in row["filename"]]
    original_total = len(original_samples)
    original_correct = sum(1 for row in original_samples if row["correct"])

    per_class = {}
    for class_name in BRICK_CLASSES:
        rows = [row for row in original_samples if row["true_label"] == class_name]
        correct = sum(1 for row in rows if row["correct"])
        per_class[class_name] = {
            "total": len(rows),
            "correct": correct,
            "accuracy": (correct / len(rows)) if rows else 0.0,
        }

    return {
        "overall_total": overall_total,
        "overall_correct": overall_correct,
        "overall_accuracy": (overall_correct / overall_total) if overall_total else 0.0,
        "original_total": original_total,
        "original_correct": original_correct,
        "original_accuracy": (original_correct / original_total) if original_total else 0.0,
        "per_class_original": per_class,
    }


def generate_comparison_chart(
    save_path: Path,
    model_names: Sequence[str],
    test_acc: Sequence[float],
    original_acc: Sequence[float],
    all_acc: Sequence[float],
) -> None:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
    plt.rcParams["axes.unicode_minus"] = False

    fig, ax = plt.subplots(figsize=(12, 6), dpi=160)
    x = list(range(len(model_names)))
    width = 0.22
    ax.bar([v - width for v in x], [v * 100 for v in test_acc], width=width, color="#C94E24", label="测试集")
    ax.bar(x, [v * 100 for v in original_acc], width=width, color="#3A6892", label="原始实拍")
    ax.bar([v + width for v in x], [v * 100 for v in all_acc], width=width, color="#6C8A35", label="全量实拍")
    ax.set_ylim(0, 100)
    ax.set_xticks(x)
    ax.set_xticklabels(model_names)
    ax.set_ylabel("准确率 (%)")
    ax.set_title("模型版本对比：测试集与真实照片准确率")
    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.legend(loc="upper left")
    for bars in ax.containers:
        ax.bar_label(bars, fmt="%.1f", padding=3, fontsize=9)
    fig.tight_layout()
    fig.savefig(save_path)
    plt.close(fig)


def generate_per_class_chart(save_path: Path, base_acc: dict, final_acc: dict) -> None:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
    plt.rcParams["axes.unicode_minus"] = False

    fig, ax = plt.subplots(figsize=(12, 6), dpi=160)
    x = list(range(len(BRICK_CLASSES)))
    width = 0.34
    base_values = [base_acc[name]["accuracy"] * 100 for name in BRICK_CLASSES]
    final_values = [final_acc[name]["accuracy"] * 100 for name in BRICK_CLASSES]
    ax.bar([v - width / 2 for v in x], base_values, width=width, color="#D8B28A", label="原始 v6")
    ax.bar([v + width / 2 for v in x], final_values, width=width, color="#C94E24", label="第二轮优化版")
    ax.set_ylim(0, 105)
    ax.set_xticks(x)
    ax.set_xticklabels(BRICK_CLASSES)
    ax.set_ylabel("原始实拍准确率 (%)")
    ax.set_title("各类别原始实拍准确率提升")
    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.legend(loc="upper left")
    for bars in ax.containers:
        ax.bar_label(bars, fmt="%.1f", padding=3, fontsize=9)
    fig.tight_layout()
    fig.savefig(save_path)
    plt.close(fig)


def set_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = TEXT,
    font_name: str = TITLE_FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float = 1.15,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    paragraph.line_spacing = line_spacing
    return box


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: Sequence[str],
    font_size: int = 20,
    color: RGBColor = TEXT,
    bullet_color: RGBColor = ACCENT,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0
        paragraph.line_spacing = 1.18
        run = paragraph.add_run()
        run.text = f"• {bullet}"
        font = run.font
        font.name = TITLE_FONT
        font.size = Pt(font_size)
        font.color.rgb = color
    return box


def add_section_header(slide, section_no: str, title: str, subtitle: str) -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(2.0), Inches(0.55)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.color.rgb = ACCENT
    add_textbox(
        slide,
        0.68,
        0.49,
        1.5,
        0.3,
        section_no,
        font_size=18,
        bold=True,
        color=SURFACE,
        font_name=EN_FONT,
    )
    add_textbox(slide, 0.55, 1.15, 5.5, 0.45, title, font_size=28, bold=True)
    add_textbox(slide, 0.55, 1.58, 6.4, 0.35, subtitle, font_size=13, color=MUTED, font_name=EN_FONT)


def add_metric_card(slide, left: float, top: float, width: float, title: str, value: str, note: str, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.4)
    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.16), Inches(1.45)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = color
    accent_bar.line.color.rgb = color
    add_textbox(slide, left + 0.25, top + 0.12, width - 0.3, 0.28, title, font_size=13, bold=True, color=MUTED)
    add_textbox(slide, left + 0.25, top + 0.45, width - 0.3, 0.42, value, font_size=24, bold=True, color=TEXT, font_name=EN_FONT)
    add_textbox(slide, left + 0.25, top + 0.94, width - 0.3, 0.28, note, font_size=11, color=MUTED)


def add_image(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_footer(slide, page_no: int) -> None:
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.18)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = LIGHT_PANEL
    slide.shapes[-1].line.color.rgb = LIGHT_PANEL
    add_textbox(slide, 11.9, 6.93, 0.6, 0.2, str(page_no), font_size=11, bold=True, color=MUTED, font_name=EN_FONT)


def create_presentation() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    base_test = load_json(ROOT / "runs" / "stage2" / "improve_raw128_autocrop_lr5e4_7class_v6" / "test_metrics_recheck_raw_20260422.json")
    round1_test = load_json(ROOT / "runs" / "stage2" / "v6_priority_train_real_1x1_2x2_2x4_20260422" / "test_metrics.json")
    round2_test = load_json(ROOT / "runs" / "stage2" / "v6_priority_train_real_round2_recover_1x2_1x4_20260422" / "test_metrics.json")
    binary_test = load_json(ROOT / "runs" / "stage2" / "square64_binary_train_20260422" / "test_metrics.json")

    base_real = compute_real_metrics(ROOT / "runs" / "stage2" / "photo_batch_eval_7class_v6_20260422" / "photo_predictions.json")
    round1_real = compute_real_metrics(ROOT / "runs" / "stage2" / "photo_batch_eval_v6_priority_train_real_1x1_2x2_2x4_20260422" / "photo_predictions.json")
    round2_real = compute_real_metrics(ROOT / "runs" / "stage2" / "photo_batch_eval_v6_priority_train_real_round2_recover_1x2_1x4_20260422" / "photo_predictions.json")
    binary_real = compute_real_metrics(ROOT / "runs" / "stage2" / "photo_batch_eval_square64_binary_20260422" / "photo_predictions.json")

    comparison_chart = ASSET_DIR / "accuracy_comparison.png"
    per_class_chart = ASSET_DIR / "per_class_comparison.png"
    generate_comparison_chart(
        comparison_chart,
        ["v6", "round1", "round2", "binary"],
        [
            float(base_test["accuracy"]),
            float(round1_test["test_accuracy"]),
            float(round2_test["test_accuracy"]),
            float(binary_test["test_accuracy"]),
        ],
        [
            base_real["original_accuracy"],
            round1_real["original_accuracy"],
            round2_real["original_accuracy"],
            binary_real["original_accuracy"],
        ],
        [
            base_real["overall_accuracy"],
            round1_real["overall_accuracy"],
            round2_real["overall_accuracy"],
            binary_real["overall_accuracy"],
        ],
    )
    generate_per_class_chart(per_class_chart, base_real["per_class_original"], round2_real["per_class_original"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.55), Inches(7.2), Inches(6.2)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = SURFACE
    hero.line.color.rgb = SURFACE
    add_textbox(slide, 0.9, 1.05, 5.8, 0.42, "人工智能应用实践", font_size=22, bold=True, color=ACCENT, font_name=EN_FONT)
    add_textbox(slide, 0.9, 1.55, 5.9, 1.0, "乐高积木图像识别项目", font_size=30, bold=True)
    add_textbox(slide, 0.9, 2.3, 5.9, 0.85, "从数据准备、CNN 训练到 Streamlit 部署，\n并围绕真实照片持续优化模型效果。", font_size=20, color=MUTED)
    add_metric_card(slide, 0.9, 3.45, 1.95, "默认模型", "round2", "当前应用已切换", ACCENT)
    add_metric_card(slide, 2.95, 3.45, 1.95, "原始实拍", "91.92%", "91 / 99", ACCENT_3)
    add_metric_card(slide, 5.0, 3.45, 1.95, "全量实拍", "83.33%", "200 / 240", ACCENT_4)
    add_bullets(
        slide,
        0.95,
        5.2,
        6.0,
        1.0,
        [
            "技术栈：PyTorch + Streamlit + Plotly",
            "任务范围：7 类 Brick 分类 + 网页演示 + 实拍优化",
        ],
        font_size=16,
    )
    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(0.55), Inches(4.7), Inches(6.2)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = LIGHT_PANEL
    right_panel.line.color.rgb = LIGHT_PANEL
    add_image(slide, ROOT / "runs" / "presentation_assets" / "sample_montage.png", 8.2, 0.9, 4.3, 5.2)
    add_textbox(slide, 8.35, 6.2, 4.0, 0.3, "样本拼图与模型优化结果概览", font_size=13, color=MUTED)
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "02", "项目概览", "PROJECT OVERVIEW")
    add_metric_card(slide, 0.7, 2.0, 2.6, "项目目标", "7 类识别", "支持真实照片分类与课堂展示", ACCENT)
    add_metric_card(slide, 3.5, 2.0, 2.6, "工程输出", "多页面应用", "训练、推理、评估链路完整", ACCENT_3)
    add_metric_card(slide, 6.3, 2.0, 2.6, "优化重点", "真实照片", "弱类补强与部署效果优先", ACCENT_4)
    add_bullets(
        slide,
        0.8,
        3.85,
        5.8,
        2.4,
        [
            "围绕课程要求完成从数据准备、模型训练到应用开发的完整项目闭环。",
            "将识别任务聚焦为 7 类 Brick：1x1、1x2、1x3、1x4、2x2、2x3、2x4。",
            "不仅关注测试集准确率，更关注真实拍照场景中的可用性和鲁棒性。",
        ],
        font_size=18,
    )
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(3.5), Inches(5.2), Inches(2.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE
    panel.line.color.rgb = LIGHT_PANEL
    add_textbox(slide, 7.5, 3.8, 4.5, 0.4, "项目交付内容", font_size=20, bold=True)
    add_bullets(
        slide,
        7.45,
        4.3,
        4.6,
        1.8,
        [
            "数据采集、预处理与 split 脚本",
            "PyTorch 训练/评估/实验对比代码",
            "Streamlit 多页面应用与推理接口",
            "报告、图表、PPT 与展示素材",
        ],
        font_size=16,
    )
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "03", "总体技术路线", "PROJECT PIPELINE")
    stages = [
        ("01 数据准备", "采集原始图片、质量检查、resize、split 划分"),
        ("02 模型训练", "基线 CNN + 超参数实验 + 最优 checkpoint 保存"),
        ("03 实拍优化", "弱类补强、热启动微调、真实照片复测"),
        ("04 应用部署", "Streamlit 多页面展示、批量/单张识别、Top-K 概率图"),
    ]
    colors = [ACCENT, ACCENT_3, ACCENT_4, ACCENT_2]
    left = 0.8
    for idx, ((title, desc), color) in enumerate(zip(stages, colors)):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + idx * 3.05), Inches(2.0), Inches(2.55), Inches(3.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = SURFACE
        box.line.color.rgb = color
        box.line.width = Pt(1.4)
        top_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left + idx * 3.05), Inches(2.0), Inches(2.55), Inches(0.22)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.color.rgb = color
        add_textbox(slide, left + idx * 3.05 + 0.2, 2.42, 2.1, 0.55, title, font_size=20, bold=True)
        add_textbox(slide, left + idx * 3.05 + 0.2, 3.15, 2.1, 1.7, desc, font_size=16, color=MUTED)
    add_textbox(
        slide,
        0.85,
        6.2,
        11.0,
        0.4,
        "关键原则：每轮训练后都回到同一批真实照片上复测，用结果驱动下一步优化。",
        font_size=17,
        bold=True,
        color=ACCENT,
    )
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "04", "数据集与识别任务", "DATASET")
    add_textbox(slide, 0.8, 2.0, 3.0, 0.4, "任务定义", font_size=22, bold=True)
    add_bullets(
        slide,
        0.8,
        2.45,
        4.8,
        2.0,
        [
            "识别对象为基础 Brick，共 7 类规格。",
            "数据来源同时包含网站采集图和真实拍照图。",
            "第二阶段以后单独关注真实照片弱类与部署效果。",
        ],
        font_size=18,
    )
    classes_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.75), Inches(4.8), Inches(1.2)
    )
    classes_box.fill.solid()
    classes_box.fill.fore_color.rgb = SURFACE
    classes_box.line.color.rgb = LIGHT_PANEL
    add_textbox(slide, 1.0, 5.0, 4.3, 0.5, "类别集合：1x1 / 1x2 / 1x3 / 1x4 / 2x2 / 2x3 / 2x4", font_size=18, bold=True)
    add_image(slide, ROOT / "runs" / "presentation_assets" / "sample_montage.png", 6.0, 1.85, 6.2, 4.85)
    add_textbox(slide, 6.15, 6.35, 5.2, 0.3, "项目样本拼图，覆盖 7 类 Brick 的多角度图像。", font_size=13, color=MUTED)
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "05", "模型与应用", "MODEL & APPLICATION")
    add_metric_card(slide, 0.8, 1.95, 2.2, "输入配置", "128×128", "auto_crop + normalize", ACCENT)
    add_metric_card(slide, 3.2, 1.95, 2.2, "模型骨干", "4 层 CNN", "baseline + GAP + FC", ACCENT_3)
    add_metric_card(slide, 5.6, 1.95, 2.2, "参数量", "876,343", "当前默认模型结构", ACCENT_4)
    add_bullets(
        slide,
        0.9,
        3.15,
        5.6,
        2.5,
        [
            "训练阶段保留成熟的 v6 结构，不盲目换网络，而是重点优化真实照片数据分布。",
            "应用阶段使用 Streamlit 多页面结构，统一走 model_def + inference 推理入口。",
            "页面支持上传识别、批量预测、拍照识别、Top-K 概率图和阶段总览。",
        ],
        font_size=17,
    )
    add_image(slide, ROOT / "assets" / "demo_outputs" / "photo_2x2_0002_detected_latest.png", 7.15, 2.2, 2.35, 3.9)
    add_image(slide, ROOT / "assets" / "demo_outputs" / "photo_1x4_0006_detected.png", 9.75, 2.2, 2.35, 3.9)
    add_textbox(slide, 7.2, 6.35, 4.8, 0.3, "检测与展示素材可直接用于课堂演示和汇报说明。", font_size=13, color=MUTED)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "06", "真实照片优化方法", "OPTIMIZATION STRATEGY")
    add_bullets(
        slide,
        0.9,
        1.95,
        5.5,
        3.8,
        [
            "第一轮：锁定弱类 1x1 / 2x2 / 2x4，在 train split 内安全复制真实照片，并从 v6 checkpoint 热启动微调。",
            "第二轮：发现 1x2 / 1x4 被副作用压低后，再定向补回 1x2 / 1x4 的 train 真实照片。",
            "始终使用小学习率热启动，而不是每轮从零训练，保留已有有效特征。",
            "避免数据泄漏：补强脚本只复制 train 内已有真实照片，不把 val/test 图塞回训练。",
        ],
        font_size=17,
    )
    timeline = [
        ("v6 基线", "测试集 87.5%\n原始实拍 65.66%"),
        ("round1", "补强 1x1/2x2/2x4\n原始实拍 78.79%"),
        ("round2", "补回 1x2/1x4\n原始实拍 91.92%"),
    ]
    positions = [7.1, 8.95, 10.8]
    for (title, desc), left_pos in zip(timeline, positions):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.3), Inches(1.55), Inches(2.65)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = ACCENT if title == "round2" else LIGHT_PANEL
        card.line.width = Pt(1.5)
        add_textbox(slide, left_pos + 0.12, 2.55, 1.25, 0.35, title, font_size=18, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left_pos + 0.1, 3.1, 1.28, 1.1, desc, font_size=14, color=MUTED, align=PP_ALIGN.CENTER)
    connector = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(8.45), Inches(3.2), Inches(0.38), Inches(0.55)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = ACCENT_3
    connector.line.color.rgb = ACCENT_3
    connector = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(10.3), Inches(3.2), Inches(0.38), Inches(0.55)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = ACCENT_3
    connector.line.color.rgb = ACCENT_3
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "07", "关键结果对比", "RESULTS")
    add_image(slide, comparison_chart, 0.8, 1.85, 8.0, 4.7)
    add_metric_card(slide, 9.2, 2.0, 3.0, "当前默认模型", "round2", "已切换到应用默认加载", ACCENT)
    add_metric_card(slide, 9.2, 3.7, 3.0, "原始实拍提升", "+26.26%", "65.66% -> 91.92%", ACCENT_3)
    add_metric_card(slide, 9.2, 5.4, 3.0, "全量实拍提升", "+31.66%", "51.67% -> 83.33%", ACCENT_4)
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "08", "分类别效果提升", "PER-CLASS IMPROVEMENT")
    add_image(slide, per_class_chart, 0.8, 1.85, 7.8, 4.8)
    add_bullets(
        slide,
        8.95,
        2.1,
        3.1,
        3.3,
        [
            "1x1：14.29% -> 85.71%",
            "1x2：91.30% -> 95.65%",
            "2x2：28.57% -> 85.71%",
            "2x4：62.07% -> 100.00%",
            "2x3：72.73% -> 90.91%",
        ],
        font_size=17,
    )
    add_textbox(
        slide,
        8.95,
        5.8,
        3.0,
        0.7,
        "说明：第二轮的核心价值，不只是整体准确率更高，而是把原来最难的弱类显著拉了起来。",
        font_size=15,
        color=MUTED,
    )
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "09", "训练曲线与部署取舍", "TRAINING EVIDENCE")
    add_image(slide, ROOT / "runs" / "stage2" / "v6_priority_train_real_round2_recover_1x2_1x4_20260422" / "history_accuracy.png", 0.8, 1.95, 5.7, 4.1)
    add_image(slide, ROOT / "runs" / "stage2" / "v6_priority_train_real_round2_recover_1x2_1x4_20260422" / "history_loss.png", 6.75, 1.95, 5.7, 4.1)
    add_metric_card(slide, 1.0, 6.2, 2.4, "best_epoch", str(round2_test["best_epoch"]), "第二轮在第 2 轮达到最佳", ACCENT)
    add_metric_card(slide, 3.7, 6.2, 2.4, "best_val_loss", f"{float(round2_test['best_val_loss']):.3f}", "验证集损失稳定", ACCENT_3)
    add_metric_card(slide, 6.4, 6.2, 2.4, "test_accuracy", f"{float(round2_test['test_accuracy']) * 100:.2f}%", "测试集略低于 v6，但更贴近真实部署目标", ACCENT_4)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "10", "RGB 与二值化对比", "RGB VS BINARY")
    add_metric_card(slide, 0.9, 2.0, 2.7, "RGB 当前默认模型", "91.92%", "原始实拍 91 / 99", ACCENT)
    add_metric_card(slide, 3.95, 2.0, 2.7, "纯二值化模型", "36.36%", "原始实拍 36 / 99", ACCENT_3)
    add_metric_card(slide, 7.0, 2.0, 2.7, "RGB 全量实拍", "83.33%", "200 / 240", ACCENT_4)
    add_metric_card(slide, 10.05, 2.0, 2.3, "二值化全量", "29.58%", "71 / 240", ACCENT_2)
    add_bullets(
        slide,
        0.95,
        4.0,
        5.6,
        2.0,
        [
            "纯二值化会丢失颜色、边界灰度和阴影信息，真实照片中容易把多类样本压到 2x4 / 2x3。",
            "因此当前线上继续使用 RGB 模型是正确路线，二值化更适合作为后续辅助特征，而不是主模型。",
        ],
        font_size=17,
    )
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(4.0), Inches(5.0), Inches(2.0)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE
    panel.line.color.rgb = LIGHT_PANEL
    add_textbox(slide, 7.25, 4.25, 4.2, 0.35, "结论", font_size=20, bold=True)
    add_textbox(
        slide,
        7.25,
        4.7,
        4.35,
        1.0,
        "当前项目应继续以 RGB 识别链路为主，\n后续如果要结合二值化，建议只作为辅助输入。",
        font_size=18,
        color=MUTED,
    )
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "11", "应用展示与落地情况", "APPLICATION")
    add_bullets(
        slide,
        0.9,
        1.95,
        5.0,
        2.3,
        [
            "项目已经落成 Streamlit 多页面应用，支持上传识别、拍照识别、批量预测、示例体验和阶段总览。",
            "默认模型已切换到第二轮优化版，应用地址为 http://127.0.0.1:8501。",
            "模型推理统一走 model_def.py + inference.py，便于后续替换 checkpoint 和继续迭代。",
        ],
        font_size=17,
    )
    add_image(slide, ROOT / "assets" / "demo_outputs" / "photo_2x2_0002_detected_latest.png", 6.3, 1.95, 2.6, 4.2)
    add_image(slide, ROOT / "assets" / "demo_outputs" / "photo_1x4_0006_detected.png", 9.2, 1.95, 2.6, 4.2)
    add_textbox(slide, 6.4, 6.35, 5.0, 0.3, "应用默认加载第二轮最优模型，已适合直接用于课堂演示。", font_size=13, color=MUTED)
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_section_header(slide, "12", "总结与下一步", "SUMMARY")
    add_bullets(
        slide,
        0.95,
        2.0,
        5.7,
        2.8,
        [
            "项目已完成：数据准备、模型训练、真实照片优化、Web 应用部署和课堂展示材料整合。",
            "当前默认模型已把原始实拍准确率提升到 91.92%，形成面向真实场景的有效闭环。",
            "剩余难点集中在 1x1 ↔ 1x2、1x4 → 2x4 等相近规格混淆。",
        ],
        font_size=18,
    )
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.05), Inches(5.15), Inches(3.2)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE
    panel.line.color.rgb = LIGHT_PANEL
    add_textbox(slide, 7.3, 2.3, 4.2, 0.35, "后续可继续优化", font_size=21, bold=True)
    add_bullets(
        slide,
        7.25,
        2.8,
        4.4,
        1.9,
        [
            "针对 hard cases 做 round3 定点补强",
            "加入轻量阴影/亮度增强",
            "为 1x1/1x2、1x4/2x4 设计二阶段复判",
        ],
        font_size=16,
    )
    add_textbox(slide, 0.95, 6.15, 11.0, 0.45, "一句话总结：这套项目已经从“能训练、能跑通”推进到“能在真实照片里稳定工作并且可展示”。", font_size=18, bold=True, color=ACCENT)
    add_footer(slide, 12)

    prs.save(OUTPUT_PPT)
    print(OUTPUT_PPT)


if __name__ == "__main__":
    create_presentation()
